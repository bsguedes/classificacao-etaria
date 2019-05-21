from pptx import Presentation
from pptx.util import Inches
from parse import parse_csv, bonus_csv
from slides import fill_slide, bonus_card

presentation = Presentation()
presentation.slide_width = Inches(2.5)
presentation.slide_height = Inches(3.5)
cards = parse_csv('cards.csv')
bonus = bonus_csv('bonus.csv')


def add_slide_from_card(c):
    layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(layout)
    fill_slide(slide, c)


def add_bonus_card(c):
    layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(layout)
    bonus_card(slide, c)


for card in cards:
    add_slide_from_card(card)


for card in bonus:
    add_bonus_card(card)


presentation.save('cards.pptx')

