import math
import os.path
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR


def text_box(slide, text, left, top, width=None, height=0.4,
             font_size=None, alignment=None, vertical_alignment=None, bold=False, italic=False, word_wrap=False):
    width = width if width is not None else 1
    tx_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx_box.text_frame
    tf.text = text
    tf.word_wrap = word_wrap
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.italic = italic
    tf.paragraphs[0].font.name = 'Arial Narrow'
    if font_size is not None:
        tf.paragraphs[0].font.size = Pt(font_size)
    if alignment is not None:
        tf.paragraphs[0].alignment = alignment
    if vertical_alignment is not None:
        tf.vertical_anchor = vertical_alignment


def add_rectangle(slide, color, left, top, width, height, border=None, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    rectangle = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    fill = rectangle.fill
    fill.solid()
    fill.fore_color.rgb = color
    line = rectangle.line
    line.color.rgb = border if border is not None else color
    rectangle.shadow.inherit = False


def add_image(slide, name, left, top, height):
    pic_path = 'img/%s.png' % name
    slide.shapes.add_picture(pic_path, Inches(left), Inches(top), height=Inches(height))


def add_photo(slide, name, left, top, height):
    pic_path = 'modphotos/%s' % name
    if os.path.isfile(pic_path):
        slide.shapes.add_picture(pic_path, Inches(left), Inches(top), height=Inches(height))


def add_programming(slide, prog):
    d = {'A1': prog[0], 'A2': prog[1], 'A3': prog[2]}
    if d['A1'] and d['A2'] and d['A3']:
        add_image(slide, 'A1', 0.11, 0.27, 0.29)
        add_image(slide, 'A2', 0.2475, 0.125, 0.29)
        add_image(slide, 'A3', 0.385, 0.27, 0.29)
    elif d['A1'] and d['A2']:
        add_image(slide, 'A1', 0.11, 0.2, 0.29)
        add_image(slide, 'A2', 0.385, 0.2, 0.29)
    elif d['A2'] and d['A3']:
        add_image(slide, 'A2', 0.11, 0.2, 0.29)
        add_image(slide, 'A3', 0.385, 0.2, 0.29)
    elif d['A1'] and d['A3']:
        add_image(slide, 'A1', 0.11, 0.2, 0.29)
        add_image(slide, 'A3', 0.385, 0.2, 0.29)
    elif d['A1']:
        add_image(slide, 'A1', 0.2475, 0.2, 0.29)
    elif d['A2']:
        add_image(slide, 'A2', 0.2475, 0.2, 0.29)
    elif d['A3']:
        add_image(slide, 'A3', 0.2475, 0.2, 0.29)


def add_tokens(slide, tokens):
    a = ['X1', 'X2', 'X3', 'X4', 'X5', 'X']
    abs_token = [math.ceil(x) for x in tokens]
    token_amount = sum(abs_token)

    if token_amount == 0:
        text_box(slide, '—', 0.1, 0.44, 0.58, font_size=14, alignment=PP_ALIGN.CENTER)
    elif token_amount == 1:
        add_image(slide, a[tokens.index(1)], 0.325, 0.59, 0.13)
    elif token_amount == 2:
        c = 0
        while sum(abs_token) > 0:
            for i in range(len(a)):
                if abs_token[i] > 0:
                    abs_token[i] -= 1
                    add_image(slide, a[i], 0.215 + c * 0.22, 0.59, 0.13)
                    c += 1
                    break
        text_box(slide, '/' if sum(tokens) < 1 else '+', 0.345, 0.52, 0.09, font_size=10, alignment=PP_ALIGN.CENTER)
    elif token_amount == 3:
        c = 0
        while sum(abs_token) > 0:
            for i in range(len(a)):
                if abs_token[i] > 0:
                    abs_token[i] -= 1
                    add_image(slide, a[i], 0.125 + c * 0.20, 0.59, 0.13)
                    c += 1
                    break
        text_box(slide, '/' if sum(tokens) < 1 else '+', 0.255, 0.52, 0.07, font_size=10, alignment=PP_ALIGN.CENTER)
        text_box(slide, '/' if sum(tokens) < 1 else '+', 0.455, 0.52, 0.07, font_size=10, alignment=PP_ALIGN.CENTER)


def fill_slide(slide, card):
    text_box(slide, card.name, 0.75, 0.1, 1.6, font_size=12, alignment=PP_ALIGN.CENTER)
    text_box(slide, card.attribute, 0.75, 0.35, 1.6, font_size=9, alignment=PP_ALIGN.CENTER, italic=True)
    add_photo(slide, card.img_path, 0.6, 0.76, 1.83)
    add_rectangle(slide, card.ability_color(), 0.0, 2.6, 2.5, 0.6)
    text_box(slide, card.ability_text(), 0.0, 2.6, 2.5, font_size=8, height=0.6, alignment=PP_ALIGN.JUSTIFY_LOW,
             vertical_alignment=MSO_ANCHOR.MIDDLE, word_wrap=True)
    # add_rectangle(slide, RGBColor(224, 224, 224), 0.1, 0.0, 0.58, 0.75, border=RGBColor(192, 192, 192))
    add_image(slide, 'star', 1.9, 0.6, 0.1)
    text_box(slide, card.year, 2.0, 0.535, font_size=8)
    text_box(slide, card.points, 0.2, 1.0, 0.15, font_size=14, alignment=PP_ALIGN.RIGHT)
    add_image(slide, 'ibope', 0.33, 1.07, 0.17)
    add_image(slide, card.age, 0.21, 1.32, 0.28)
    text_box(slide, card.money, 0.15, 1.55, 0.4, font_size=12, word_wrap=True)
    text_box(slide, card.lore, 0.0, 3.165, 2, height=0.3, font_size=6.5, alignment=PP_ALIGN.JUSTIFY_LOW,
             vertical_alignment=MSO_ANCHOR.MIDDLE, italic=True, word_wrap=True)
    add_programming(slide, card.programming)
    add_tokens(slide, card.tokens)


def bonus_card(slide, card):
    add_rectangle(slide, RGBColor(180, 216, 231), 0.4, 0.4, 1.7, 0.6)
    text_box(slide, card.name, 0.4, 0.4, 1.7, 0.6, font_size=11, alignment=PP_ALIGN.CENTER,
             vertical_alignment=MSO_ANCHOR.MIDDLE, word_wrap=True)
    text_box(slide, card.percent_text(), 0.2, 2.9, 2.1, font_size=8, alignment=PP_ALIGN.CENTER, italic=True)
    text_box(slide, card.text, 0.3, 1.8, 1.9, font_size=9, height=0.6, alignment=PP_ALIGN.CENTER,
             vertical_alignment=MSO_ANCHOR.MIDDLE, word_wrap=True)
    if card.type() == 'range':
        text_box(slide, card.p1, 0.65, 1.25, 0.15, font_size=14, alignment=PP_ALIGN.RIGHT)
        add_image(slide, 'ibope', 0.78, 1.32, 0.17)
        text_box(slide, "%s atrações" % card.range1, 0.35, 1.55, 0.9, font_size=8, alignment=PP_ALIGN.CENTER)
        text_box(slide, card.p2, 1.6, 1.25, 0.15, font_size=14, alignment=PP_ALIGN.RIGHT)
        add_image(slide, 'ibope', 1.73, 1.32, 0.17)
        text_box(slide, "%s atrações" % card.range2, 1.3, 1.55, 0.9, font_size=8, alignment=PP_ALIGN.CENTER)
    else:
        text_box(slide, card.score, 1.1, 1.25, 0.15, font_size=14, alignment=PP_ALIGN.RIGHT)
        add_image(slide, 'ibope', 1.23, 1.32, 0.17)
        text_box(slide, 'por atração', 0.8, 1.55, 0.9, font_size=8, alignment=PP_ALIGN.CENTER)


